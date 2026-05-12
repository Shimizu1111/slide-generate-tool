"""デザイン解析モジュール。

画像・pptx・PDFからデザイン要素（レイアウト、配色、フォント、余白、装飾など）を
解析し、テンプレート構築に使える構造化データとして返す。

- pptx: python-pptx で直接解析
- 画像: Claude Vision API でデザインを解析
- PDF: pdf2image で画像化してから Vision API で解析
"""

from __future__ import annotations

import base64
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ============================================================
#  画像からデザイン解析（Claude Vision API）
# ============================================================

def analyze_image(image_path: str | Path) -> dict[str, Any]:
    """画像ファイルからスライドデザインを解析する。

    Claude Vision API を使って画像を解析し、
    レイアウト・配色・フォント・装飾などの構造化データを返す。
    """
    import anthropic

    image_path = Path(image_path)
    if not image_path.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")

    image_data = base64.standard_b64encode(image_path.read_bytes()).decode("utf-8")

    suffix = image_path.suffix.lower()
    media_type_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
    }
    media_type = media_type_map.get(suffix, "image/png")

    client = anthropic.Anthropic()

    system_prompt = """あなたはスライドデザインの解析エキスパートです。
アップロードされた画像を分析し、このデザインをpython-pptxで再現するための
構造化データをJSON形式で返してください。

以下の情報を抽出してください:
1. スライドサイズ（アスペクト比から推定: 16:9=254x143mm, 4:3=254x190.5mm）
2. 背景（色、グラデーション、画像の有無）
3. 各要素の詳細:
   - 種類（テキスト、図形、画像、アイコン）
   - 位置（left_mm, top_mm）と大きさ（width_mm, height_mm）
   - テキスト内容、フォントサイズ(pt)、太字/斜体、色(RGB hex)
   - 図形の種類、塗り色、線色
   - 配置（左寄せ/中央/右寄せ）

レスポンスは以下のJSON形式のみで返してください（説明文不要）:
```json
{
  "slide_width_mm": 254.0,
  "slide_height_mm": 143.0,
  "background": {
    "type": "solid",
    "color": "FFFFFF"
  },
  "elements": [
    {
      "type": "textbox",
      "left_mm": 25.0,
      "top_mm": 20.0,
      "width_mm": 200.0,
      "height_mm": 30.0,
      "text": "タイトルテキスト",
      "font_name": "Arial",
      "font_size_pt": 36,
      "font_color": "333333",
      "bold": true,
      "italic": false,
      "alignment": "left"
    },
    {
      "type": "shape",
      "shape_type": "rectangle",
      "left_mm": 10.0,
      "top_mm": 130.0,
      "width_mm": 234.0,
      "height_mm": 3.0,
      "fill_color": "0066CC",
      "line_color": null
    }
  ]
}
```"""

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_type,
                        "data": image_data,
                    },
                },
                {
                    "type": "text",
                    "text": "この画像のスライドデザインを解析してください。",
                },
            ],
        }],
        system=system_prompt,
    )

    response_text = response.content[0].text

    # JSONを抽出
    if "```json" in response_text:
        json_str = response_text.split("```json")[1].split("```")[0].strip()
    elif "```" in response_text:
        json_str = response_text.split("```")[1].split("```")[0].strip()
    else:
        json_str = response_text

    return json.loads(json_str)


def analyze_pdf(pdf_path: str | Path, page: int = 0) -> dict[str, Any]:
    """PDFファイルからスライドデザインを解析する。

    PDFを画像に変換してからVision APIで解析する。
    """
    import tempfile
    from pdf2image import convert_from_path

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    images = convert_from_path(str(pdf_path), first_page=page + 1, last_page=page + 1, dpi=200)
    if not images:
        raise ValueError(f"Failed to convert PDF page {page}")

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        images[0].save(tmp.name, "PNG")
        return analyze_image(tmp.name)


# ============================================================
#  pptx からデザイン解析（python-pptx で直接解析）
# ============================================================

def analyze_pptx(pptx_path: str | Path) -> dict[str, Any]:
    """pptxファイルの全構造を解析して返す。"""
    prs = Presentation(str(pptx_path))
    return {
        "slide_width_mm": emu_to_mm(prs.slide_width),
        "slide_height_mm": emu_to_mm(prs.slide_height),
        "slide_layouts": [_inspect_layout(layout) for layout in prs.slide_layouts],
        "slides": [_inspect_slide(slide) for slide in prs.slides],
    }


def analyze_pptx_summary(pptx_path: str | Path) -> dict[str, Any]:
    """pptxの概要をコンパクトに返す（API送信用）。"""
    prs = Presentation(str(pptx_path))
    layouts = []
    for layout in prs.slide_layouts:
        phs = []
        for ph in layout.placeholders:
            phs.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "left_mm": emu_to_mm(ph.left),
                "top_mm": emu_to_mm(ph.top),
                "width_mm": emu_to_mm(ph.width),
                "height_mm": emu_to_mm(ph.height),
            })
        layouts.append({"name": layout.name, "placeholders": phs})

    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            s = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_mm": emu_to_mm(shape.left),
                "top_mm": emu_to_mm(shape.top),
                "width_mm": emu_to_mm(shape.width),
                "height_mm": emu_to_mm(shape.height),
            }
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text]
                if texts:
                    s["texts"] = texts
                for p in shape.text_frame.paragraphs:
                    if p.runs:
                        font = p.runs[0].font
                        if font.name:
                            s["font"] = font.name
                        if font.size:
                            s["font_size_pt"] = font.size.pt
                        if font.bold:
                            s["bold"] = True
                        try:
                            if font.color and font.color.type is not None:
                                s["color"] = str(font.color.rgb)
                        except AttributeError:
                            pass
                        break
            shapes.append(s)
        slides.append({
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": shapes,
        })

    return {
        "slide_width_mm": emu_to_mm(prs.slide_width),
        "slide_height_mm": emu_to_mm(prs.slide_height),
        "layouts": layouts,
        "slides": slides,
    }


# ============================================================
#  汎用エントリポイント（ファイル種別を自動判定）
# ============================================================

def analyze(file_path: str | Path) -> dict[str, Any]:
    """ファイル種別を自動判定してデザインを解析する。"""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        return analyze_pptx(file_path)
    elif suffix == ".pdf":
        return analyze_pdf(file_path)
    elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        return analyze_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


# ============================================================
#  内部ヘルパー（pptx解析用）
# ============================================================

def _inspect_layout(layout) -> dict[str, Any]:
    return {
        "name": layout.name,
        "placeholders": [_inspect_placeholder(ph) for ph in layout.placeholders],
    }


def _inspect_slide(slide) -> dict[str, Any]:
    layout_name = slide.slide_layout.name if slide.slide_layout else None
    return {
        "layout_name": layout_name,
        "shapes": [_inspect_shape(shape) for shape in slide.shapes],
    }


def _inspect_shape(shape) -> dict[str, Any]:
    info: dict[str, Any] = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": emu_to_mm(shape.left),
        "top": emu_to_mm(shape.top),
        "width": emu_to_mm(shape.width),
        "height": emu_to_mm(shape.height),
    }

    if shape.has_text_frame:
        info["text_frame"] = _inspect_text_frame(shape.text_frame)

    try:
        pf = shape.placeholder_format
        if pf is not None:
            info["placeholder_idx"] = pf.idx
            info["placeholder_type"] = str(pf.type)
    except ValueError:
        pass

    if shape.has_table:
        info["table"] = _inspect_table(shape.table)

    if hasattr(shape, "image") and shape.shape_type == 13:
        info["image"] = {
            "content_type": shape.image.content_type,
            "width": emu_to_mm(shape.width),
            "height": emu_to_mm(shape.height),
        }

    return info


def _inspect_placeholder(ph) -> dict[str, Any]:
    info: dict[str, Any] = {
        "idx": ph.placeholder_format.idx,
        "type": str(ph.placeholder_format.type),
        "name": ph.name,
        "left": emu_to_mm(ph.left),
        "top": emu_to_mm(ph.top),
        "width": emu_to_mm(ph.width),
        "height": emu_to_mm(ph.height),
    }
    if ph.has_text_frame:
        info["text_frame"] = _inspect_text_frame(ph.text_frame)
    return info


def _inspect_text_frame(text_frame) -> dict[str, Any]:
    return {
        "word_wrap": text_frame.word_wrap,
        "paragraphs": [_inspect_paragraph(p) for p in text_frame.paragraphs],
    }


def _inspect_paragraph(paragraph) -> dict[str, Any]:
    info: dict[str, Any] = {
        "text": paragraph.text,
        "alignment": _alignment_str(paragraph.alignment),
        "level": paragraph.level,
    }
    if paragraph.space_before is not None:
        info["space_before"] = str(paragraph.space_before)
    if paragraph.space_after is not None:
        info["space_after"] = str(paragraph.space_after)
    if paragraph.line_spacing is not None:
        info["line_spacing"] = str(paragraph.line_spacing)
    if paragraph.runs:
        info["runs"] = [_inspect_run(run) for run in paragraph.runs]
    return info


def _inspect_run(run) -> dict[str, Any]:
    font = run.font
    info: dict[str, Any] = {"text": run.text}
    if font.name is not None:
        info["font_name"] = font.name
    if font.size is not None:
        info["font_size_pt"] = font.size.pt
    if font.bold is not None:
        info["bold"] = font.bold
    if font.italic is not None:
        info["italic"] = font.italic
    if font.underline is not None:
        info["underline"] = font.underline
    try:
        if font.color and font.color.type is not None:
            info["color_rgb"] = str(font.color.rgb)
    except AttributeError:
        pass
    return info


def _inspect_table(table) -> dict[str, Any]:
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append({
                "text": cell.text,
                "text_frame": _inspect_text_frame(cell.text_frame),
            })
        rows.append(cells)
    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "data": rows,
    }


def _alignment_str(alignment) -> str | None:
    if alignment is None:
        return None
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "distribute",
    }
    return mapping.get(alignment, str(alignment))


def emu_to_mm(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / 914400 * 25.4, 2)
