"""テンプレート解析モジュール。

pptxファイルのレイアウト・プレースホルダー・フォント・色・サイズなどを
解析して構造化データとして返す。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def inspect_template(pptx_path: str | Path) -> dict[str, Any]:
    """テンプレートの全構造を解析して返す。"""
    prs = Presentation(str(pptx_path))
    return {
        "slide_width": emu_to_mm(prs.slide_width),
        "slide_height": emu_to_mm(prs.slide_height),
        "slide_layouts": [_inspect_layout(layout) for layout in prs.slide_layouts],
        "slides": [_inspect_slide(slide) for slide in prs.slides],
    }


def inspect_slide(pptx_path: str | Path, slide_index: int = 0) -> dict[str, Any]:
    """特定スライドの詳細を返す。"""
    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    if slide_index >= len(slides):
        raise IndexError(f"Slide index {slide_index} out of range (total: {len(slides)})")
    return _inspect_slide(slides[slide_index])


def _inspect_layout(layout) -> dict[str, Any]:
    """スライドレイアウトを解析。"""
    return {
        "name": layout.name,
        "placeholders": [_inspect_placeholder(ph) for ph in layout.placeholders],
    }


def _inspect_slide(slide) -> dict[str, Any]:
    """スライドを解析。"""
    layout_name = slide.slide_layout.name if slide.slide_layout else None
    return {
        "layout_name": layout_name,
        "shapes": [_inspect_shape(shape) for shape in slide.shapes],
    }


def _inspect_shape(shape) -> dict[str, Any]:
    """シェイプを解析。"""
    info: dict[str, Any] = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": emu_to_mm(shape.left),
        "top": emu_to_mm(shape.top),
        "width": emu_to_mm(shape.width),
        "height": emu_to_mm(shape.height),
    }

    if shape.has_text_frame:
        info["text_frame"] = _inspect_text_frame(shape.text_frame)

    try:
        pf = shape.placeholder_format
        if pf is not None:
            info["placeholder_idx"] = pf.idx
            info["placeholder_type"] = str(pf.type)
    except ValueError:
        pass

    if shape.has_table:
        info["table"] = _inspect_table(shape.table)

    if hasattr(shape, "image") and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        info["image"] = {
            "content_type": shape.image.content_type,
            "width": emu_to_mm(shape.width),
            "height": emu_to_mm(shape.height),
        }

    return info


def _inspect_placeholder(ph) -> dict[str, Any]:
    """プレースホルダーを解析。"""
    info: dict[str, Any] = {
        "idx": ph.placeholder_format.idx,
        "type": str(ph.placeholder_format.type),
        "name": ph.name,
        "left": emu_to_mm(ph.left),
        "top": emu_to_mm(ph.top),
        "width": emu_to_mm(ph.width),
        "height": emu_to_mm(ph.height),
    }

    if ph.has_text_frame:
        info["text_frame"] = _inspect_text_frame(ph.text_frame)

    return info


def _inspect_text_frame(text_frame) -> dict[str, Any]:
    """テキストフレームを解析。"""
    return {
        "word_wrap": text_frame.word_wrap,
        "paragraphs": [_inspect_paragraph(p) for p in text_frame.paragraphs],
    }


def _inspect_paragraph(paragraph) -> dict[str, Any]:
    """パラグラフを解析。"""
    info: dict[str, Any] = {
        "text": paragraph.text,
        "alignment": _alignment_str(paragraph.alignment),
        "level": paragraph.level,
    }

    if paragraph.space_before is not None:
        info["space_before"] = str(paragraph.space_before)
    if paragraph.space_after is not None:
        info["space_after"] = str(paragraph.space_after)
    if paragraph.line_spacing is not None:
        info["line_spacing"] = str(paragraph.line_spacing)

    if paragraph.runs:
        info["runs"] = [_inspect_run(run) for run in paragraph.runs]

    return info


def _inspect_run(run) -> dict[str, Any]:
    """ランを解析。"""
    font = run.font
    info: dict[str, Any] = {
        "text": run.text,
    }

    if font.name is not None:
        info["font_name"] = font.name
    if font.size is not None:
        info["font_size_pt"] = font.size.pt
    if font.bold is not None:
        info["bold"] = font.bold
    if font.italic is not None:
        info["italic"] = font.italic
    if font.underline is not None:
        info["underline"] = font.underline
    try:
        if font.color and font.color.type is not None:
            info["color_rgb"] = str(font.color.rgb)
    except AttributeError:
        pass

    return info


def _inspect_table(table) -> dict[str, Any]:
    """テーブルを解析。"""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append({
                "text": cell.text,
                "text_frame": _inspect_text_frame(cell.text_frame),
            })
        rows.append(cells)
    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "data": rows,
    }


def _alignment_str(alignment) -> str | None:
    """アライメントを文字列に変換。"""
    if alignment is None:
        return None
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "distribute",
    }
    return mapping.get(alignment, str(alignment))


def inspect_template_summary(pptx_path: str | Path) -> dict[str, Any]:
    """テンプレートの概要をコンパクトに返す（API送信用）。"""
    prs = Presentation(str(pptx_path))
    layouts = []
    for layout in prs.slide_layouts:
        phs = []
        for ph in layout.placeholders:
            phs.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "left_mm": emu_to_mm(ph.left),
                "top_mm": emu_to_mm(ph.top),
                "width_mm": emu_to_mm(ph.width),
                "height_mm": emu_to_mm(ph.height),
            })
        layouts.append({"name": layout.name, "placeholders": phs})

    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            s = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_mm": emu_to_mm(shape.left),
                "top_mm": emu_to_mm(shape.top),
                "width_mm": emu_to_mm(shape.width),
                "height_mm": emu_to_mm(shape.height),
            }
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text]
                if texts:
                    s["texts"] = texts
                # 最初のランのフォント情報だけ取得
                for p in shape.text_frame.paragraphs:
                    if p.runs:
                        font = p.runs[0].font
                        if font.name:
                            s["font"] = font.name
                        if font.size:
                            s["font_size_pt"] = font.size.pt
                        if font.bold:
                            s["bold"] = True
                        try:
                            if font.color and font.color.type is not None:
                                s["color"] = str(font.color.rgb)
                        except AttributeError:
                            pass
                        break
            shapes.append(s)
        slides.append({
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": shapes,
        })

    return {
        "slide_width_mm": emu_to_mm(prs.slide_width),
        "slide_height_mm": emu_to_mm(prs.slide_height),
        "layouts": layouts,
        "slides": slides,
    }


def emu_to_mm(emu: int | None) -> float | None:
    """EMU値をmm単位に変換。"""
    if emu is None:
        return None
    return round(emu / 914400 * 25.4, 2)
