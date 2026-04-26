"""スライド生成エンジン。

テンプレート + 入力データ → 完成スライドを生成する。
コードはプレースホルダーへの流し込みのみ行い、デザインには一切触れない。
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


def generate_slides(
    template_path: str | Path,
    input_data: dict[str, Any] | str | Path,
    output_path: str | Path,
) -> Path:
    """テンプレートと入力データからスライドを生成する。

    Args:
        template_path: テンプレート.pptxのパス
        input_data: 入力データ（dict、JSONファイルパス、またはJSON文字列）
        output_path: 出力先パス

    Returns:
        生成されたファイルのパス
    """
    template_path = Path(template_path)
    output_path = Path(output_path)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    data = _load_input_data(input_data)
    prs = Presentation(str(template_path))

    slides_data = data.get("slides", [])
    for slide_data in slides_data:
        _add_slide(prs, slide_data)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _load_input_data(input_data: dict[str, Any] | str | Path) -> dict[str, Any]:
    """入力データを読み込んでdictにする。"""
    if isinstance(input_data, dict):
        return input_data

    path = Path(str(input_data))
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)

    # JSON文字列として解析
    if isinstance(input_data, str):
        return json.loads(input_data)

    raise ValueError(f"Cannot load input data: {input_data}")


def _add_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    """スライドを1枚追加して内容を流し込む。"""
    layout_index = slide_data.get("layout", 0)
    layout_name = slide_data.get("layout_name")

    # レイアウト名で検索（指定されている場合）
    layout = None
    if layout_name:
        for sl in prs.slide_layouts:
            if sl.name == layout_name:
                layout = sl
                break
        if layout is None:
            raise ValueError(
                f"Layout '{layout_name}' not found. "
                f"Available: {[sl.name for sl in prs.slide_layouts]}"
            )
    else:
        if layout_index >= len(prs.slide_layouts):
            raise IndexError(
                f"Layout index {layout_index} out of range "
                f"(total: {len(prs.slide_layouts)})"
            )
        layout = prs.slide_layouts[layout_index]

    slide = prs.slides.add_slide(layout)

    # プレースホルダーに内容を流し込む
    placeholders = slide_data.get("placeholders", {})
    for key, value in placeholders.items():
        idx = int(key) if isinstance(key, str) and key.isdigit() else key
        if idx in slide.placeholders:
            ph = slide.placeholders[idx]
            if isinstance(value, str):
                ph.text = value
            elif isinstance(value, dict):
                _fill_placeholder(ph, value)

    # 画像プレースホルダー
    images = slide_data.get("images", {})
    for key, image_path in images.items():
        idx = int(key) if isinstance(key, str) and key.isdigit() else key
        if idx in slide.placeholders:
            ph = slide.placeholders[idx]
            ph.insert_picture(str(image_path))


def _fill_placeholder(placeholder, content: dict[str, Any]) -> None:
    """プレースホルダーに構造化データを流し込む。

    テンプレートのデフォルト書式を継承する。
    フォントやサイズは変更しない。
    """
    if "text" in content:
        placeholder.text = content["text"]

    if "paragraphs" in content:
        tf = placeholder.text_frame
        tf.clear()
        for i, para_data in enumerate(content["paragraphs"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            text = para_data.get("text", "")
            level = para_data.get("level", 0)
            p.text = text
            p.level = level


def validate_input_data(
    template_path: str | Path,
    input_data: dict[str, Any] | str | Path,
) -> list[str]:
    """入力データがテンプレートと整合するか検証する。

    Returns:
        問題点のリスト（空ならOK）
    """
    template_path = Path(template_path)
    data = _load_input_data(input_data)
    prs = Presentation(str(template_path))

    issues: list[str] = []
    available_layouts = {sl.name: i for i, sl in enumerate(prs.slide_layouts)}

    for i, slide_data in enumerate(data.get("slides", [])):
        layout_name = slide_data.get("layout_name")
        layout_index = slide_data.get("layout", 0)

        if layout_name and layout_name not in available_layouts:
            issues.append(
                f"Slide {i}: layout '{layout_name}' not found in template. "
                f"Available: {list(available_layouts.keys())}"
            )
            continue

        # レイアウトのプレースホルダーと入力データのキーを比較
        if layout_name:
            idx = available_layouts[layout_name]
        else:
            idx = layout_index

        if idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[idx]
            available_phs = {ph.placeholder_format.idx for ph in layout.placeholders}
            requested_phs = set()
            for key in slide_data.get("placeholders", {}):
                requested_phs.add(int(key) if isinstance(key, str) and key.isdigit() else key)

            missing = requested_phs - available_phs
            if missing:
                issues.append(
                    f"Slide {i}: placeholder indices {missing} not found in layout. "
                    f"Available: {available_phs}"
                )

    return issues
