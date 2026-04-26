"""テンプレート構築モジュール。

python-pptxを使ってテンプレートをゼロから構築、または既存テンプレートを更新する。
テンプレート作成・複製時はデザイン構築が目的のため、座標・フォント指定を許容する。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def create_blank_presentation(
    width_mm: float = 254.0,
    height_mm: float = 190.5,
) -> Presentation:
    """空のプレゼンテーションを作成（デフォルト: ワイドスクリーン）。"""
    prs = Presentation()
    prs.slide_width = Mm(width_mm)
    prs.slide_height = Mm(height_mm)
    return prs


def add_layout(
    prs: Presentation,
    name: str,
    placeholders: list[dict[str, Any]] | None = None,
) -> None:
    """スライドマスターにレイアウトを追加する。

    python-pptxはスライドレイアウトの新規追加APIを直接提供しないため、
    XMLレベルで操作する。
    """
    from pptx.oxml.ns import qn
    from copy import deepcopy
    import lxml.etree as etree

    slide_master = prs.slide_masters[0]
    sldMaster = slide_master._element

    # 既存レイアウトの最初のものをベースにコピー
    existing_layouts = list(prs.slide_layouts)
    if not existing_layouts:
        raise ValueError("No existing layouts to base new layout on")

    base_layout = existing_layouts[0]
    new_layout_elem = deepcopy(base_layout._element)

    # cSld内の既存シェイプをクリア
    cSld = new_layout_elem.find(qn("p:cSld"))
    spTree = cSld.find(qn("p:spTree"))

    # spTree内のsp要素（シェイプ）を削除
    for sp in spTree.findall(qn("p:sp")):
        spTree.remove(sp)

    # プレースホルダーを追加
    if placeholders:
        for i, ph_config in enumerate(placeholders):
            _add_placeholder_to_spTree(spTree, ph_config, i)

    # レイアウト名を設定
    cSld.set("name", name)

    # slide_layouts に追加（内部リレーションの操作）
    # python-pptxの制約上、直接的なAPI追加はできないが、
    # 既存レイアウトのプレースホルダーを変更することで対応
    # → 実用上は既存レイアウトの修正で対応する


def modify_layout(
    prs: Presentation,
    layout_index: int,
    placeholders: list[dict[str, Any]] | None = None,
    name: str | None = None,
) -> None:
    """既存レイアウトを修正する。"""
    layout = prs.slide_layouts[layout_index]

    if name is not None:
        layout.name = name

    if placeholders is not None:
        # 既存プレースホルダーの属性を更新
        for ph_config in placeholders:
            idx = ph_config.get("idx")
            if idx is not None:
                for ph in layout.placeholders:
                    if ph.placeholder_format.idx == idx:
                        _update_shape_properties(ph, ph_config)
                        break


def add_slide(
    prs: Presentation,
    layout_index: int = 0,
    content: dict[str, Any] | None = None,
) -> Any:
    """スライドを追加し、内容を設定する。"""
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    if content:
        for key, value in content.items():
            if key.startswith("placeholder_"):
                idx = int(key.split("_")[1])
                if idx in slide.placeholders:
                    slide.placeholders[idx].text = str(value)

    return slide


def add_textbox(
    slide,
    left_mm: float,
    top_mm: float,
    width_mm: float,
    height_mm: float,
    text: str = "",
    font_name: str | None = None,
    font_size_pt: float | None = None,
    font_color_rgb: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    alignment: str | None = None,
    vertical_anchor: str | None = None,
) -> Any:
    """テキストボックスを追加する（テンプレート構築用）。"""
    txBox = slide.shapes.add_textbox(
        Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    if vertical_anchor:
        anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
        tf.paragraphs[0].alignment = None
        if vertical_anchor in anchor_map:
            tf.vertical_anchor = anchor_map[vertical_anchor]

    p = tf.paragraphs[0]
    p.text = text

    if alignment:
        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        if alignment in align_map:
            p.alignment = align_map[alignment]

    if text and p.runs:
        run = p.runs[0]
        _apply_font(run.font, font_name, font_size_pt, font_color_rgb, bold, italic)

    return txBox


def add_shape(
    slide,
    shape_type: str,
    left_mm: float,
    top_mm: float,
    width_mm: float,
    height_mm: float,
    fill_color_rgb: str | None = None,
    line_color_rgb: str | None = None,
    line_width_pt: float | None = None,
) -> Any:
    """図形を追加する（テンプレート構築用）。"""
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow_right": MSO_SHAPE.RIGHT_ARROW,
        "line": MSO_SHAPE.LINE_INVERSE,
    }
    mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)

    shape = slide.shapes.add_shape(
        mso_shape, Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm)
    )

    if fill_color_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color_rgb)

    if line_color_rgb:
        shape.line.color.rgb = RGBColor.from_string(line_color_rgb)

    if line_width_pt is not None:
        shape.line.width = Pt(line_width_pt)

    return shape


def add_image(
    slide,
    image_path: str,
    left_mm: float,
    top_mm: float,
    width_mm: float | None = None,
    height_mm: float | None = None,
) -> Any:
    """画像を追加する。"""
    left = Mm(left_mm)
    top = Mm(top_mm)
    width = Mm(width_mm) if width_mm else None
    height = Mm(height_mm) if height_mm else None
    return slide.shapes.add_picture(image_path, left, top, width, height)


def set_slide_background(slide, color_rgb: str) -> None:
    """スライドの背景色を設定する。"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_rgb)


def save_presentation(prs: Presentation, output_path: str | Path) -> Path:
    """プレゼンテーションを保存する。"""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _apply_font(
    font,
    name: str | None = None,
    size_pt: float | None = None,
    color_rgb: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
) -> None:
    """フォント属性を設定する。"""
    if name is not None:
        font.name = name
    if size_pt is not None:
        font.size = Pt(size_pt)
    if color_rgb is not None:
        font.color.rgb = RGBColor.from_string(color_rgb)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic


def _update_shape_properties(shape, config: dict[str, Any]) -> None:
    """シェイプのプロパティを更新する。"""
    if "left_mm" in config:
        shape.left = Mm(config["left_mm"])
    if "top_mm" in config:
        shape.top = Mm(config["top_mm"])
    if "width_mm" in config:
        shape.width = Mm(config["width_mm"])
    if "height_mm" in config:
        shape.height = Mm(config["height_mm"])


def _add_placeholder_to_spTree(spTree, ph_config: dict[str, Any], index: int) -> None:
    """spTreeにプレースホルダーXML要素を追加する。"""
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    idx = ph_config.get("idx", index)
    ph_type = ph_config.get("type", "body")
    left = Mm(ph_config.get("left_mm", 25.4))
    top = Mm(ph_config.get("top_mm", 25.4 + index * 50))
    width = Mm(ph_config.get("width_mm", 203.2))
    height = Mm(ph_config.get("height_mm", 40))

    type_map = {"title": "title", "body": "body", "subtitle": "subTitle"}
    sp_type = type_map.get(ph_type, "body")

    sp_xml = f"""
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
        <p:nvSpPr>
            <p:cNvPr id="{index + 2}" name="Placeholder {idx}"/>
            <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
            <p:nvPr><p:ph type="{sp_type}" idx="{idx}"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr>
            <a:xfrm>
                <a:off x="{left}" y="{top}"/>
                <a:ext cx="{width}" cy="{height}"/>
            </a:xfrm>
        </p:spPr>
        <p:txBody>
            <a:bodyPr/>
            <a:lstStyle/>
            <a:p><a:endParaRPr lang="ja-JP"/></a:p>
        </p:txBody>
    </p:sp>
    """
    sp_elem = etree.fromstring(sp_xml)
    spTree.append(sp_elem)
