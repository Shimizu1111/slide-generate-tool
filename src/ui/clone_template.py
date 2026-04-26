"""テンプレート複製ページ。

既存のpptxのデザインを完全に保持したまま、
中身をクリアして再利用可能なテンプレートに変換する。

左ペイン: 元テンプレートのプレビュー
右ペイン: テンプレート化後のプレビュー
下部: 差分表示と改善チャット
"""

from __future__ import annotations

import json
import shutil
import tempfile
from pathlib import Path

import streamlit as st

from src.core.inspector import inspect_template, inspect_template_summary
from src.core.comparator import compare_templates
from src.core.renderer import render_slides, check_libreoffice
from src.core.builder import save_presentation


def render_page():
    st.title("テンプレート複製")
    st.caption("既存のpptxをデザインそのままでテンプレート化します")

    # セッション状態の初期化
    if "clone_source_path" not in st.session_state:
        st.session_state.clone_source_path = None
    if "clone_target_path" not in st.session_state:
        st.session_state.clone_target_path = None
    if "clone_source_images" not in st.session_state:
        st.session_state.clone_source_images = []
    if "clone_target_images" not in st.session_state:
        st.session_state.clone_target_images = []
    if "clone_diffs" not in st.session_state:
        st.session_state.clone_diffs = None
    if "clone_messages" not in st.session_state:
        st.session_state.clone_messages = []

    if "clone_reference_image" not in st.session_state:
        st.session_state.clone_reference_image = None

    # ── サイドバー ──
    with st.sidebar:
        st.subheader("テンプレートファイル")
        source_file = st.file_uploader("元の .pptx をアップロード", type=["pptx"], key="clone_source")
        if source_file:
            source_path = Path(tempfile.mkdtemp(prefix="clone_src_")) / source_file.name
            source_path.write_bytes(source_file.getvalue())
            st.session_state.clone_source_path = str(source_path)
            _update_source_preview()

        ref_image = st.file_uploader(
            "理想のデザイン画像 (任意)",
            type=["png", "jpg", "jpeg"],
            key="clone_ref_img",
            help="Canvaやブラウザのスクリーンショットなど、本来こう見えるべきという画像",
        )
        if ref_image:
            ref_path = Path(tempfile.mkdtemp(prefix="clone_ref_")) / ref_image.name
            ref_path.write_bytes(ref_image.getvalue())
            st.session_state.clone_reference_image = str(ref_path)

        if st.session_state.clone_source_path and not st.session_state.clone_target_path:
            if st.button("テンプレート化を開始"):
                _start_cloning()

        if st.session_state.clone_target_path:
            pptx_path = Path(st.session_state.clone_target_path)
            if pptx_path.exists():
                st.divider()
                st.subheader("テンプレートを保存")

                template_name = st.text_input(
                    "テンプレート名",
                    value="my_template",
                    key="clone_template_name",
                )
                if st.button("templates/ に保存"):
                    _save_template(pptx_path, template_name)

                with open(pptx_path, "rb") as f:
                    st.download_button(
                        label="ダウンロード",
                        data=f.read(),
                        file_name=f"{template_name}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

    # ── メインエリア ──
    if st.session_state.clone_source_path:
        # リファレンス画像がある場合は3カラム、なければ2カラム
        has_ref = st.session_state.clone_reference_image is not None

        if has_ref:
            ref_col, left_col, right_col = st.columns([1, 1, 1])
            with ref_col:
                st.subheader("理想のデザイン")
                st.image(st.session_state.clone_reference_image, use_container_width=True)
                st.caption("この見た目に合わせます")
        else:
            left_col, right_col = st.columns([1, 1])

        with left_col:
            st.subheader("現在のpptx")
            _show_preview(st.session_state.clone_source_images, "source_slide_idx")

        with right_col:
            st.subheader("テンプレート化後")
            if st.session_state.clone_target_images:
                _show_preview(st.session_state.clone_target_images, "target_slide_idx")
            else:
                st.info("「テンプレート化を開始」を押してください")

        # ── 差分表示 ──
        if st.session_state.clone_diffs is not None:
            st.divider()
            _show_diff_report()

        # ── 改善チャット ──
        if st.session_state.clone_target_path:
            st.divider()
            _show_refinement_chat()


def _show_preview(images: list[Path], slider_key: str):
    """スライドプレビューを表示。"""
    if not images:
        st.info("プレビューなし")
        return
    total = len(images)
    if total > 1:
        idx = st.slider("スライド", 0, total - 1, 0, key=slider_key)
    else:
        idx = 0
    st.image(str(images[idx]), use_container_width=True)


def _show_diff_report():
    """差分レポートを表示。"""
    diffs = st.session_state.clone_diffs
    total = diffs["total_diffs"]

    if total == 0:
        st.success("元テンプレートと完全一致しています。")
        return

    st.info(f"差分: {total}件（テキストをクリアした分など）")

    summary = diffs.get("summary", {})
    if summary:
        cols = st.columns(len(summary))
        label_map = {
            "font": "フォント", "color": "色", "position_size": "位置・サイズ",
            "alignment": "配置", "layout": "レイアウト", "other": "その他",
        }
        for i, (category, count) in enumerate(summary.items()):
            with cols[i % len(cols)]:
                st.metric(label_map.get(category, category), count)

    with st.expander("差分詳細"):
        for diff in diffs["diffs"]:
            if diff["type"] == "value_mismatch":
                unit = diff.get("unit", "")
                st.text(f"  {diff['path']}: {diff['source']}{unit} → {diff['target']}{unit}")
            else:
                st.text(f"  {diff['path']}: {diff['type']}")


def _show_refinement_chat():
    """改善用チャットを表示。"""
    st.subheader("テンプレートを調整")

    # 対象スライド選択
    target_images = st.session_state.clone_target_images
    total = len(target_images) if target_images else 0

    if total > 0:
        options = ["すべてのスライド"] + [f"{i + 1}枚目" for i in range(total)]
        selected = st.multiselect(
            "修正対象のスライド",
            options,
            default=["すべてのスライド"],
            key="clone_slide_target",
        )

        if not selected or "すべてのスライド" in selected:
            target = "all"
        else:
            indices = []
            for s in selected:
                if s != "すべてのスライド":
                    idx = int(s.replace("枚目", "")) - 1
                    indices.append(idx)
            target = ",".join(str(i) for i in sorted(indices))
    else:
        target = "all"

    if st.session_state.clone_reference_image:
        st.caption("理想のデザイン画像と見比べて、違いを指示してください")
    else:
        st.caption("デザインの微調整をチャットで指示できます")

    for msg in st.session_state.clone_messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("調整指示を入力（例: タイトルのフォントをもう少し大きく）"):
        # 対象スライド情報を付加
        if target == "all":
            target_label = "すべてのスライド"
        else:
            indices = [int(i) for i in target.split(",")]
            target_label = "、".join(f"{i+1}枚目" for i in indices)

        display_msg = f"【対象: {target_label}】{prompt}"
        st.session_state.clone_messages.append({"role": "user", "content": display_msg})
        with st.chat_message("user"):
            st.markdown(display_msg)

        with st.chat_message("assistant"):
            with st.spinner("修正中..."):
                response = _process_refinement(prompt, target)
                st.markdown(response)

        st.session_state.clone_messages.append({"role": "assistant", "content": response})
        st.rerun()


def _save_template(pptx_path: Path, name: str):
    """テンプレートを templates/ ディレクトリに保存する。"""
    templates_dir = Path(__file__).resolve().parent.parent.parent / "templates"
    templates_dir.mkdir(exist_ok=True)
    dest = templates_dir / f"{name}.pptx"
    shutil.copy2(str(pptx_path), str(dest))
    st.success(f"保存しました: templates/{name}.pptx")
    st.code(f"make generate T=templates/{name}.pptx I=input/data.json", language="bash")


# ── 複製ロジック ──

def _start_cloning():
    """テンプレートの複製を開始する。"""
    source_path = st.session_state.clone_source_path
    if not source_path:
        return

    with st.spinner("テンプレート化中..."):
        target_path = _clone_direct(source_path)
        if target_path:
            st.session_state.clone_target_path = target_path
            _update_target_preview()
            _update_diffs()
            st.rerun()


def _clone_direct(source_path: str) -> str | None:
    """元pptxを直接コピーしてテンプレート化する。

    デザイン（背景、図形、テーマ、フォント、色、レイアウト）は完全に保持。
    プレースホルダーのテキストのみクリアして、再利用可能なテンプレートにする。
    """
    from pptx import Presentation

    output_dir = Path(tempfile.mkdtemp(prefix="clone_tgt_"))
    output_path = output_dir / "cloned_template.pptx"

    # ファイルをそのままコピー（全XML・メディア・テーマを保持）
    shutil.copy2(source_path, str(output_path))

    # プレースホルダーのテキストをクリア（書式は保持）
    prs = Presentation(str(output_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                ph = shape.placeholder_format
                if ph is not None and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
            except ValueError:
                # placeholder_formatがないシェイプはスキップ
                continue

    prs.save(str(output_path))
    return str(output_path)


def _process_refinement(prompt: str, target_slides: str = "all") -> str:
    """ユーザーの調整指示を処理する。"""
    from openai import OpenAI

    client = OpenAI()

    target_info = inspect_template_summary(st.session_state.clone_target_path)

    # 対象スライドの説明
    if target_slides == "all":
        slide_instruction = "すべてのスライドに対して修正してください。"
        slide_code_hint = "for slide in prs.slides:"
    else:
        indices = [int(i) for i in target_slides.split(",")]
        labels = ", ".join(str(i) for i in indices)
        slide_instruction = f"スライドインデックス {labels} (0始まり) のみ修正してください。他のスライドは変更しないでください。"
        slide_code_hint = f"target_indices = [{labels}]\nfor i in target_indices:\n    slide = prs.slides[i]"

    system_prompt = f"""あなたはpptxテンプレートを調整するエキスパートです。
ユーザーの指示に基づいて、テンプレートを修正するpython-pptxコードを生成してください。

{slide_instruction}

対象スライドへのアクセス例:
```python
{slide_code_hint}
```

レスポンスは以下のJSON形式で返してください:
```json
{{
  "message": "ユーザーへの返答",
  "code": "修正用Pythonコード"
}}
```

変数 `prs` に現在のPresentationオブジェクトが入っています。
修正後も `prs` に残してください。save()は呼ばないでください。

利用可能なインポート:
- from pptx import Presentation
- from pptx.util import Inches, Pt, Mm, Emu
- from pptx.dml.color import RGBColor
- from pptx.enum.text import PP_ALIGN, MSO_ANCHOR"""

    ref_note = ""
    if st.session_state.clone_reference_image:
        ref_note = "\n\n注意: ユーザーは理想のデザイン画像を参照しながら差分を指摘しています。位置やサイズの修正指示には正確に対応してください。"

    messages = [{
        "role": "user",
        "content": (
            f"現在のテンプレート構造:\n```json\n{json.dumps(target_info, ensure_ascii=False, indent=2)}\n```"
            f"{ref_note}\n\n"
            f"ユーザーの指示: {prompt}"
        ),
    }]

    response = client.chat.completions.create(
        model="gpt-5.4",
        max_completion_tokens=4096,
        messages=[{"role": "system", "content": system_prompt}] + messages,
    )

    response_text = response.choices[0].message.content

    try:
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        else:
            json_str = response_text
        result = json.loads(json_str)
        message = result.get("message", "修正しました。")
        code = result.get("code")
    except (json.JSONDecodeError, IndexError):
        return response_text

    if code:
        new_path = _execute_refinement_code(code)
        if new_path:
            st.session_state.clone_target_path = new_path
            _update_target_preview()
            _update_diffs()

    return message


def _execute_refinement_code(code: str) -> str | None:
    """修正コードを実行する。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Mm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation(st.session_state.clone_target_path)

    namespace = {
        "prs": prs,
        "Presentation": Presentation,
        "Inches": Inches, "Pt": Pt, "Mm": Mm, "Emu": Emu,
        "RGBColor": RGBColor,
        "PP_ALIGN": PP_ALIGN, "MSO_ANCHOR": MSO_ANCHOR,
    }

    try:
        exec(code, namespace)
    except Exception as e:
        st.error(f"修正コード実行エラー: {e}")
        return None

    prs = namespace.get("prs")
    if prs is None:
        return None

    output_path = Path(tempfile.mkdtemp(prefix="clone_refined_")) / "cloned_template.pptx"
    save_presentation(prs, output_path)
    return str(output_path)


# ── プレビュー更新 ──

def _update_source_preview():
    path = st.session_state.clone_source_path
    if path and Path(path).exists() and check_libreoffice():
        try:
            st.session_state.clone_source_images = render_slides(path)
        except Exception:
            st.session_state.clone_source_images = []


def _update_target_preview():
    path = st.session_state.clone_target_path
    if path and Path(path).exists() and check_libreoffice():
        try:
            st.session_state.clone_target_images = render_slides(path)
        except Exception:
            st.session_state.clone_target_images = []


def _update_diffs():
    source = st.session_state.clone_source_path
    target = st.session_state.clone_target_path
    if source and target:
        try:
            st.session_state.clone_diffs = compare_templates(source, target)
        except Exception as e:
            st.error(f"差分比較エラー: {e}")
            st.session_state.clone_diffs = None
