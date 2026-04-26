"""スライド生成ページ。

テンプレート内の既存スライドを選び、各スライドの内容を指定して
完成スライドを生成する。
"""

from __future__ import annotations

import json
import tempfile
import uuid
from pathlib import Path

import streamlit as st
from PIL import Image

from src.core.inspector import inspect_template_summary
from src.core.generator import generate_slides, validate_input_data
from src.core.renderer import render_slides, check_libreoffice


_TEMPLATES_DIR = Path(__file__).resolve().parent.parent.parent / "templates"

# フッター・日付・ページ番号のプレースホルダーは非表示
_SKIP_PH_TYPES = {"DATE", "FOOTER", "SLIDE_NUMBER"}
_SKIP_PH_IDX = {10, 11, 12}  # 一般的なフッター系IDX


def render_page():
    st.title("スライド生成")
    st.caption("テンプレートのスライドを選んで、内容を入れ替えます")

    # セッション状態の初期化
    for key, default in {
        "gen_template_path": None,
        "gen_template_info": None,
        "gen_slide_previews": None,
        "gen_slides": [],
        "gen_preview_images": None,
        "gen_output_path": None,
        "gen_refine_messages": [],
    }.items():
        if key not in st.session_state:
            st.session_state[key] = default

    # ── サイドバー ──
    with st.sidebar:
        st.subheader("テンプレート")
        _template_selector()

        if st.session_state.gen_slides:
            st.divider()
            st.subheader("生成")
            if st.button("プレビュー & 生成", type="primary"):
                _generate()

            if st.session_state.gen_output_path:
                output_path = Path(st.session_state.gen_output_path)
                if output_path.exists():
                    with open(output_path, "rb") as f:
                        st.download_button(
                            label="ダウンロード (.pptx)",
                            data=f.read(),
                            file_name="generated.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                    save_name = st.text_input("ファイル名", value="generated", key="gen_save_name")
                    if st.button("output/ に保存"):
                        import shutil
                        output_dir = Path(__file__).resolve().parent.parent.parent / "output"
                        output_dir.mkdir(exist_ok=True)
                        dest = output_dir / f"{save_name}.pptx"
                        shutil.copy2(str(output_path), str(dest))
                        st.success(f"保存: output/{save_name}.pptx")

    # ── メインエリア ──
    if not st.session_state.gen_template_path:
        st.info("サイドバーからテンプレートを選択してください")
        return

    # スライドテンプレートギャラリー
    _slide_gallery()

    st.divider()

    # スライド一覧の編集
    if st.session_state.gen_slides:
        _slide_editor()
    else:
        st.info("上のスライドをクリックして追加してください")

    # プレビュー + 修正チャット
    if st.session_state.gen_preview_images:
        st.divider()
        _show_preview()
        st.divider()
        _show_refinement_chat()


def _template_selector():
    """テンプレート選択UI。"""
    _TEMPLATES_DIR.mkdir(exist_ok=True)
    templates = sorted(_TEMPLATES_DIR.glob("*.pptx"))

    if not templates:
        st.warning("templates/ にテンプレートがありません")
        st.caption("「テンプレート作成」か「テンプレート複製」で作成してください")
        return

    names = [t.stem for t in templates]
    selected = st.selectbox("テンプレートを選択", names, key="gen_template_select")

    if selected:
        path = _TEMPLATES_DIR / f"{selected}.pptx"
        if str(path) != st.session_state.gen_template_path:
            st.session_state.gen_template_path = str(path)
            st.session_state.gen_template_info = inspect_template_summary(path)
            st.session_state.gen_slides = []
            st.session_state.gen_preview_images = None
            st.session_state.gen_output_path = None
            # テンプレート内のスライドをプレビュー
            if check_libreoffice():
                with st.spinner("スライドプレビュー生成中..."):
                    st.session_state.gen_slide_previews = render_slides(path, dpi=100)
            else:
                st.session_state.gen_slide_previews = None


def _slide_gallery():
    """テンプレート内のスライドをギャラリー表示。クリックで追加。"""
    info = st.session_state.gen_template_info
    previews = st.session_state.gen_slide_previews
    if not info or not info.get("slides"):
        st.warning("テンプレートにスライドがありません")
        return

    st.subheader("スライドを選んで追加")
    st.caption("使いたいスライドの「追加」を押すと、下の構成に追加されます")

    slides = info["slides"]
    cols_per_row = 4
    for row_start in range(0, len(slides), cols_per_row):
        cols = st.columns(cols_per_row)
        for i, col in enumerate(cols):
            idx = row_start + i
            if idx >= len(slides):
                break
            slide = slides[idx]
            with col:
                # プレビュー画像
                if previews and idx < len(previews):
                    st.image(str(previews[idx]), use_container_width=True)

                layout_name = slide.get("layout") or f"スライド {idx + 1}"
                st.caption(f"{idx + 1}. {layout_name}")
                if st.button("追加", key=f"add_slide_{idx}"):
                    _add_slide_from_template(idx)
                    st.rerun()


def _add_slide_from_template(slide_index: int):
    """テンプレートのスライドをベースに追加する。"""
    info = st.session_state.gen_template_info
    slide = info["slides"][slide_index]

    # シェイプからプレースホルダー的なテキスト入力フィールドを抽出
    placeholders = {}
    for shape in slide.get("shapes", []):
        # テキストを持つシェイプだけ
        if "texts" not in shape and "font" not in shape:
            continue

        name = shape.get("name", "")
        # フッター・ページ番号系はスキップ
        if any(skip in name.lower() for skip in ["footer", "slide number", "date"]):
            continue

        texts = shape.get("texts", [])
        current_text = "\n".join(texts) if texts else ""

        placeholders[name] = {
            "label": _shape_label(shape),
            "value": current_text,
            "multiline": len(texts) > 1 or shape.get("font_size_pt", 100) < 20,
        }

    st.session_state.gen_slides.append({
        "id": str(uuid.uuid4()),
        "slide_index": slide_index,
        "layout_name": slide.get("layout") or f"スライド {slide_index + 1}",
        "placeholders": placeholders,
    })


def _shape_label(shape: dict) -> str:
    """シェイプから分かりやすいラベルを生成。"""
    name = shape.get("name", "")
    texts = shape.get("texts", [])
    font_size = shape.get("font_size_pt")

    # フォントサイズから推定
    if font_size and font_size >= 24:
        return "タイトル"
    if font_size and font_size >= 18:
        return "見出し"

    if "title" in name.lower():
        return "タイトル"
    if "subtitle" in name.lower():
        return "サブタイトル"
    if "content" in name.lower() or "body" in name.lower():
        return "本文"

    if texts:
        # 最初のテキストの一部をラベルに
        preview = texts[0][:15]
        return f"テキスト ({preview}...)" if len(texts[0]) > 15 else f"テキスト ({preview})"

    return "テキスト"


def _slide_editor():
    """スライド一覧の編集UI。"""
    st.subheader(f"スライド構成 ({len(st.session_state.gen_slides)}枚)")

    slides = st.session_state.gen_slides
    previews = st.session_state.gen_slide_previews
    to_delete = None
    to_move = None

    for i, slide in enumerate(slides):
        with st.container(border=True):
            # ヘッダー行
            header_cols = st.columns([3, 1, 1, 1, 1])
            with header_cols[0]:
                st.markdown(f"**{i + 1}枚目** — {slide['layout_name']}")
            with header_cols[1]:
                if i > 0 and st.button("↑", key=f"up_{slide['id']}"):
                    to_move = (i, i - 1)
            with header_cols[2]:
                if i < len(slides) - 1 and st.button("↓", key=f"down_{slide['id']}"):
                    to_move = (i, i + 1)
            with header_cols[3]:
                if st.button("複製", key=f"dup_{slide['id']}"):
                    import copy
                    new_slide = copy.deepcopy(slide)
                    new_slide["id"] = str(uuid.uuid4())
                    slides.insert(i + 1, new_slide)
                    st.rerun()
            with header_cols[4]:
                if st.button("削除", key=f"del_{slide['id']}"):
                    to_delete = i

            # プレビューと入力フォームを横並び
            preview_col, form_col = st.columns([1, 2])

            with preview_col:
                slide_idx = slide.get("slide_index", 0)
                if previews and slide_idx < len(previews):
                    st.image(str(previews[slide_idx]), use_container_width=True)

            with form_col:
                for ph_name, ph_data in slide["placeholders"].items():
                    label = ph_data["label"]
                    if ph_data.get("multiline"):
                        new_val = st.text_area(
                            label,
                            value=ph_data["value"],
                            height=80,
                            key=f"ph_{slide['id']}_{ph_name}",
                        )
                    else:
                        new_val = st.text_input(
                            label,
                            value=ph_data["value"],
                            key=f"ph_{slide['id']}_{ph_name}",
                        )
                    ph_data["value"] = new_val

    if to_delete is not None:
        slides.pop(to_delete)
        st.rerun()
    if to_move is not None:
        a, b = to_move
        slides[a], slides[b] = slides[b], slides[a]
        st.rerun()


def _generate():
    """テンプレートをベースにスライドを生成する。"""
    from pptx import Presentation
    import shutil

    template_path = st.session_state.gen_template_path
    slides_config = st.session_state.gen_slides
    info = st.session_state.gen_template_info

    if not slides_config:
        st.warning("スライドが追加されていません")
        return

    with st.spinner("スライドを生成中..."):
        output_path = Path(tempfile.mkdtemp(prefix="gen_output_")) / "generated.pptx"

        # テンプレートをコピー
        shutil.copy2(template_path, str(output_path))
        prs = Presentation(str(output_path))

        # 元スライドの数
        original_count = len(prs.slides)

        # 選択されたスライドをコピーして内容を差し替え
        # まず必要なスライドを追加
        from pptx.oxml.ns import qn
        from copy import deepcopy
        import lxml.etree as etree

        new_slides_xml = []
        for slide_config in slides_config:
            src_idx = slide_config["slide_index"]
            if src_idx >= original_count:
                continue

            # 元スライドのXMLをコピー
            src_slide = prs.slides[src_idx]
            new_slide_elem = deepcopy(src_slide._element)
            new_slides_xml.append((new_slide_elem, slide_config, src_slide))

        # 新しいプレゼンテーションを構築（元テンプレートベース）
        # 選択されたスライドのみ含むpptxを作る
        result_prs = Presentation(template_path)

        # 全既存スライドを取得
        existing_slides = list(result_prs.slides)

        # 既存スライドを全削除
        for _ in range(len(existing_slides)):
            rId = result_prs.slides._sldIdLst[0].get(qn("r:id"))
            result_prs.part.drop_rel(rId)
            result_prs.slides._sldIdLst.remove(result_prs.slides._sldIdLst[0])

        # 選択されたスライドを順に追加（元テンプレートから再度読み込み）
        for slide_config in slides_config:
            src_idx = slide_config["slide_index"]

            # 元テンプレートを都度読み込んでスライドをコピー
            tmp_prs = Presentation(template_path)
            if src_idx >= len(tmp_prs.slides):
                continue

            src_slide = tmp_prs.slides[src_idx]
            layout = src_slide.slide_layout

            # 同名レイアウトをresult_prsから見つける
            target_layout = None
            for tl in result_prs.slide_layouts:
                if tl.name == layout.name:
                    target_layout = tl
                    break
            if target_layout is None:
                target_layout = result_prs.slide_layouts[0]

            new_slide = result_prs.slides.add_slide(target_layout)

            # 元スライドのシェイプを新スライドにコピー
            for shape in src_slide.shapes:
                el = deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            # テキストを差し替え（空欄ならテキストをクリア）
            ph_values = slide_config.get("placeholders", {})
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_name = shape.name
                if shape_name in ph_values:
                    new_text = ph_values[shape_name]["value"]
                    if new_text:
                        _replace_text_preserve_format(shape.text_frame, new_text)
                    else:
                        # 空欄 → 書式を保持しつつテキストをクリア
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = ""

        result_prs.save(str(output_path))
        st.session_state.gen_output_path = str(output_path)

        # プレビュー
        if check_libreoffice():
            try:
                images = render_slides(output_path)
                st.session_state.gen_preview_images = images
            except Exception:
                st.session_state.gen_preview_images = None

    st.rerun()


def _replace_text_preserve_format(text_frame, new_text: str):
    """テキストフレームのテキストを差し替え（書式は保持）。"""
    lines = new_text.split("\n")

    for i, paragraph in enumerate(text_frame.paragraphs):
        if i >= len(lines):
            # 余分なパラグラフはクリア
            for run in paragraph.runs:
                run.text = ""
            continue

        if paragraph.runs:
            # 最初のランにテキストを入れ、残りはクリア
            paragraph.runs[0].text = lines[i]
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = lines[i]

    # 元のパラグラフ数より行が多い場合は追加
    if len(lines) > len(text_frame.paragraphs):
        for line in lines[len(text_frame.paragraphs):]:
            p = text_frame.add_paragraph()
            p.text = line


def _show_preview():
    """生成結果のプレビュー。"""
    st.subheader("生成結果プレビュー")
    images = st.session_state.gen_preview_images
    if not images:
        return

    total = len(images)
    if total > 1:
        idx = st.slider("スライド", 0, total - 1, 0, key="gen_preview_idx")
    else:
        idx = 0
    st.image(str(images[idx]), use_container_width=True)


def _slide_target_selector(key_prefix: str, total_slides: int) -> str:
    """修正対象スライドの選択UI。返り値は選択内容の説明文字列。"""
    options = ["すべてのスライド"] + [f"{i + 1}枚目" for i in range(total_slides)]
    selected = st.multiselect(
        "修正対象のスライド",
        options,
        default=["すべてのスライド"],
        key=f"{key_prefix}_slide_target",
    )

    if not selected or "すべてのスライド" in selected:
        return "all"

    # "3枚目" → 2 (0-indexed)
    indices = []
    for s in selected:
        if s != "すべてのスライド":
            idx = int(s.replace("枚目", "")) - 1
            indices.append(idx)
    return ",".join(str(i) for i in sorted(indices))


def _show_refinement_chat():
    """生成結果の修正チャット。"""
    st.subheader("デザインを修正")

    images = st.session_state.gen_preview_images or []
    total = len(images)

    if total > 0:
        target = _slide_target_selector("gen", total)
    else:
        target = "all"

    st.caption("プレビューを見て気になる点を指示してください")

    for msg in st.session_state.gen_refine_messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("修正指示を入力", key="gen_refine_input"):
        # 対象スライド情報を付加
        if target == "all":
            target_label = "すべてのスライド"
        else:
            indices = [int(i) for i in target.split(",")]
            target_label = "、".join(f"{i+1}枚目" for i in indices)

        display_msg = f"【対象: {target_label}】{prompt}"
        st.session_state.gen_refine_messages.append({"role": "user", "content": display_msg})
        with st.chat_message("user"):
            st.markdown(display_msg)

        with st.chat_message("assistant"):
            with st.spinner("修正中..."):
                response = _process_refinement(prompt, target)
                st.markdown(response)

        st.session_state.gen_refine_messages.append({"role": "assistant", "content": response})
        st.rerun()


def _process_refinement(prompt: str, target_slides: str = "all") -> str:
    """修正指示を処理する。"""
    from openai import OpenAI

    client = OpenAI()

    output_path = st.session_state.gen_output_path
    if not output_path:
        return "生成されたスライドがありません。先に生成してください。"

    target_info = inspect_template_summary(output_path)

    # 対象スライドの説明
    if target_slides == "all":
        slide_instruction = "すべてのスライドに対して修正してください。"
        slide_code_hint = "for slide in prs.slides:"
    else:
        indices = [int(i) for i in target_slides.split(",")]
        labels = ", ".join(str(i) for i in indices)
        slide_instruction = f"スライドインデックス {labels} (0始まり) のみ修正してください。他のスライドは変更しないでください。"
        slide_code_hint = f"target_indices = [{labels}]\nfor i in target_indices:\n    slide = prs.slides[i]"

    system_prompt = f"""あなたはpptxスライドのデザインを修正するエキスパートです。
ユーザーが生成されたスライドの見た目の問題を指摘します。
python-pptxを使って修正するコードを生成してください。

{slide_instruction}

対象スライドへのアクセス例:
```python
{slide_code_hint}
```

レスポンスは以下のJSON形式で返してください:
```json
{{
  "message": "ユーザーへの返答（何を修正したか）",
  "code": "修正用Pythonコード"
}}
```

変数 `prs` に現在のPresentationオブジェクトが入っています。
修正後も `prs` に残してください。save()は呼ばないでください。

利用可能なインポート:
- from pptx import Presentation
- from pptx.util import Inches, Pt, Mm, Emu
- from pptx.dml.color import RGBColor
- from pptx.enum.text import PP_ALIGN, MSO_ANCHOR"""

    messages = [{
        "role": "user",
        "content": (
            f"現在のスライド構造:\n```json\n{json.dumps(target_info, ensure_ascii=False, indent=2)}\n```\n\n"
            f"修正指示: {prompt}"
        ),
    }]

    response = client.chat.completions.create(
        model="gpt-5.4",
        max_completion_tokens=4096,
        messages=[{"role": "system", "content": system_prompt}] + messages,
    )

    response_text = response.choices[0].message.content

    try:
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        else:
            json_str = response_text
        result = json.loads(json_str)
        message = result.get("message", "修正しました。")
        code = result.get("code")
    except (json.JSONDecodeError, IndexError):
        return response_text

    if code:
        new_path = _execute_refinement_code(code)
        if new_path:
            st.session_state.gen_output_path = new_path
            if check_libreoffice():
                try:
                    images = render_slides(new_path)
                    st.session_state.gen_preview_images = images
                except Exception:
                    pass

    return message


def _execute_refinement_code(code: str) -> str | None:
    """修正コードを実行する。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Mm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation(st.session_state.gen_output_path)

    namespace = {
        "prs": prs,
        "Presentation": Presentation,
        "Inches": Inches, "Pt": Pt, "Mm": Mm, "Emu": Emu,
        "RGBColor": RGBColor,
        "PP_ALIGN": PP_ALIGN, "MSO_ANCHOR": MSO_ANCHOR,
    }

    try:
        exec(code, namespace)
    except Exception as e:
        st.error(f"修正コード実行エラー: {e}")
        return None

    prs = namespace.get("prs")
    if prs is None:
        return None

    output_path = Path(tempfile.mkdtemp(prefix="gen_refined_")) / "generated.pptx"
    prs.save(str(output_path))
    return str(output_path)
