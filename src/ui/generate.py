"""スライド生成ページ。

テンプレートを選んで、チャットまたはフォームでコンテンツを入力し、
完成スライド(.pptx)を生成する。
"""

from __future__ import annotations

import json
import tempfile
import uuid
from pathlib import Path

import streamlit as st
from PIL import Image

from src.core.analyzer import analyze_pptx_summary
from src.core.generator import generate_slides, validate_input_data
from src.core.renderer import render_slides, check_libreoffice


_TEMPLATES_DIR = Path(__file__).resolve().parent.parent.parent / "templates"

_SKIP_PH_TYPES = {"DATE", "FOOTER", "SLIDE_NUMBER"}
_SKIP_PH_IDX = {10, 11, 12}


def render_page():
    st.title("スライド生成")
    st.caption("テンプレートのデザインで自分のコンテンツのスライドを作ります")

    # セッション状態の初期化
    for key, default in {
        "gen_template_path": None,
        "gen_template_info": None,
        "gen_slide_previews": None,
        "gen_slides": [],
        "gen_preview_images": None,
        "gen_output_path": None,
        "gen_chat_messages": [],
        "gen_mode": "form",
    }.items():
        if key not in st.session_state:
            st.session_state[key] = default

    # ── サイドバー ──
    with st.sidebar:
        st.subheader("テンプレート")
        _template_selector()

        if st.session_state.gen_template_path:
            st.divider()
            st.subheader("入力方式")
            mode = st.radio(
                "コンテンツの入力方式",
                ["フォーム", "チャット"],
                key="gen_mode_radio",
                horizontal=True,
            )
            st.session_state.gen_mode = "form" if mode == "フォーム" else "chat"

        if st.session_state.gen_mode == "form" and st.session_state.gen_slides:
            st.divider()
            st.subheader("生成")
            if st.button("プレビュー & 生成", type="primary"):
                _generate()

        if st.session_state.gen_output_path:
            output_path = Path(st.session_state.gen_output_path)
            if output_path.exists():
                with open(output_path, "rb") as f:
                    st.download_button(
                        label="ダウンロード (.pptx)",
                        data=f.read(),
                        file_name="generated.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                save_name = st.text_input("ファイル名", value="generated", key="gen_save_name")
                if st.button("output/ に保存"):
                    import shutil
                    output_dir = Path(__file__).resolve().parent.parent.parent / "output"
                    output_dir.mkdir(exist_ok=True)
                    dest = output_dir / f"{save_name}.pptx"
                    shutil.copy2(str(output_path), str(dest))
                    st.success(f"保存: output/{save_name}.pptx")

    # ── メインエリア ──
    if not st.session_state.gen_template_path:
        st.info("サイドバーからテンプレートを選択してください")
        return

    if st.session_state.gen_mode == "form":
        _render_form_mode()
    else:
        _render_chat_mode()


def _template_selector():
    """テンプレート選択UI。"""
    _TEMPLATES_DIR.mkdir(exist_ok=True)
    templates = sorted(_TEMPLATES_DIR.glob("*.pptx"))

    if not templates:
        st.warning("templates/ にテンプレートがありません")
        st.caption("「デザインキャプチャ」でテンプレートを作成してください")
        return

    names = [t.stem for t in templates]
    selected = st.selectbox("テンプレートを選択", names, key="gen_template_select")

    if selected:
        path = _TEMPLATES_DIR / f"{selected}.pptx"
        if str(path) != st.session_state.gen_template_path:
            st.session_state.gen_template_path = str(path)
            st.session_state.gen_template_info = analyze_pptx_summary(path)
            st.session_state.gen_slides = []
            st.session_state.gen_preview_images = None
            st.session_state.gen_output_path = None
            st.session_state.gen_chat_messages = []
            if check_libreoffice():
                with st.spinner("スライドプレビュー生成中..."):
                    st.session_state.gen_slide_previews = render_slides(path, dpi=100)
            else:
                st.session_state.gen_slide_previews = None


# ============================================================
#  フォームモード
# ============================================================

def _render_form_mode():
    """フォームでコンテンツを入力するモード。"""
    _slide_gallery()
    st.divider()

    if st.session_state.gen_slides:
        _slide_editor()
    else:
        st.info("上のスライドをクリックして追加してください")

    if st.session_state.gen_preview_images:
        st.divider()
        _show_preview()


def _slide_gallery():
    """テンプレート内のスライドをギャラリー表示。"""
    info = st.session_state.gen_template_info
    previews = st.session_state.gen_slide_previews
    if not info or not info.get("slides"):
        st.warning("テンプレートにスライドがありません")
        return

    st.subheader("スライドを選んで追加")
    st.caption("使いたいスライドの「追加」を押すと、下の構成に追加されます")

    slides = info["slides"]
    cols_per_row = 4
    for row_start in range(0, len(slides), cols_per_row):
        cols = st.columns(cols_per_row)
        for i, col in enumerate(cols):
            idx = row_start + i
            if idx >= len(slides):
                break
            slide = slides[idx]
            with col:
                if previews and idx < len(previews):
                    st.image(str(previews[idx]), use_container_width=True)

                layout_name = slide.get("layout") or f"スライド {idx + 1}"
                st.caption(f"{idx + 1}. {layout_name}")
                if st.button("追加", key=f"add_slide_{idx}"):
                    _add_slide_from_template(idx)
                    st.rerun()


def _add_slide_from_template(slide_index: int):
    """テンプレートのスライドをベースに追加する。"""
    info = st.session_state.gen_template_info
    slide = info["slides"][slide_index]

    placeholders = {}
    for shape in slide.get("shapes", []):
        if "texts" not in shape and "font" not in shape:
            continue

        name = shape.get("name", "")
        if any(skip in name.lower() for skip in ["footer", "slide number", "date"]):
            continue

        texts = shape.get("texts", [])
        current_text = "\n".join(texts) if texts else ""

        placeholders[name] = {
            "label": _shape_label(shape),
            "value": current_text,
            "multiline": len(texts) > 1 or shape.get("font_size_pt", 100) < 20,
        }

    st.session_state.gen_slides.append({
        "id": str(uuid.uuid4()),
        "slide_index": slide_index,
        "layout_name": slide.get("layout") or f"スライド {slide_index + 1}",
        "placeholders": placeholders,
    })


def _shape_label(shape: dict) -> str:
    name = shape.get("name", "")
    texts = shape.get("texts", [])
    font_size = shape.get("font_size_pt")

    if font_size and font_size >= 24:
        return "タイトル"
    if font_size and font_size >= 18:
        return "見出し"
    if "title" in name.lower():
        return "タイトル"
    if "subtitle" in name.lower():
        return "サブタイトル"
    if "content" in name.lower() or "body" in name.lower():
        return "本文"
    if texts:
        preview = texts[0][:15]
        return f"テキスト ({preview}...)" if len(texts[0]) > 15 else f"テキスト ({preview})"
    return "テキスト"


def _slide_editor():
    """スライド一覧の編集UI。"""
    st.subheader(f"スライド構成 ({len(st.session_state.gen_slides)}枚)")

    slides = st.session_state.gen_slides
    previews = st.session_state.gen_slide_previews
    to_delete = None
    to_move = None

    for i, slide in enumerate(slides):
        with st.container(border=True):
            header_cols = st.columns([3, 1, 1, 1, 1])
            with header_cols[0]:
                st.markdown(f"**{i + 1}枚目** — {slide['layout_name']}")
            with header_cols[1]:
                if i > 0 and st.button("↑", key=f"up_{slide['id']}"):
                    to_move = (i, i - 1)
            with header_cols[2]:
                if i < len(slides) - 1 and st.button("↓", key=f"down_{slide['id']}"):
                    to_move = (i, i + 1)
            with header_cols[3]:
                if st.button("複製", key=f"dup_{slide['id']}"):
                    import copy
                    new_slide = copy.deepcopy(slide)
                    new_slide["id"] = str(uuid.uuid4())
                    slides.insert(i + 1, new_slide)
                    st.rerun()
            with header_cols[4]:
                if st.button("削除", key=f"del_{slide['id']}"):
                    to_delete = i

            preview_col, form_col = st.columns([1, 2])

            with preview_col:
                slide_idx = slide.get("slide_index", 0)
                if previews and slide_idx < len(previews):
                    st.image(str(previews[slide_idx]), use_container_width=True)

            with form_col:
                for ph_name, ph_data in slide["placeholders"].items():
                    label = ph_data["label"]
                    if ph_data.get("multiline"):
                        new_val = st.text_area(
                            label, value=ph_data["value"], height=80,
                            key=f"ph_{slide['id']}_{ph_name}",
                        )
                    else:
                        new_val = st.text_input(
                            label, value=ph_data["value"],
                            key=f"ph_{slide['id']}_{ph_name}",
                        )
                    ph_data["value"] = new_val

    if to_delete is not None:
        slides.pop(to_delete)
        st.rerun()
    if to_move is not None:
        a, b = to_move
        slides[a], slides[b] = slides[b], slides[a]
        st.rerun()


def _generate():
    """テンプレートをベースにスライドを生成する。"""
    from pptx import Presentation
    import shutil
    from pptx.oxml.ns import qn
    from copy import deepcopy

    template_path = st.session_state.gen_template_path
    slides_config = st.session_state.gen_slides

    if not slides_config:
        st.warning("スライドが追加されていません")
        return

    with st.spinner("スライドを生成中..."):
        output_path = Path(tempfile.mkdtemp(prefix="gen_output_")) / "generated.pptx"
        shutil.copy2(template_path, str(output_path))

        result_prs = Presentation(str(output_path))

        # 既存スライドを全削除
        existing_slides = list(result_prs.slides)
        for _ in range(len(existing_slides)):
            rId = result_prs.slides._sldIdLst[0].get(qn("r:id"))
            result_prs.part.drop_rel(rId)
            result_prs.slides._sldIdLst.remove(result_prs.slides._sldIdLst[0])

        # 選択されたスライドを順に追加
        for slide_config in slides_config:
            src_idx = slide_config["slide_index"]

            tmp_prs = Presentation(template_path)
            if src_idx >= len(tmp_prs.slides):
                continue

            src_slide = tmp_prs.slides[src_idx]
            layout = src_slide.slide_layout

            target_layout = None
            for tl in result_prs.slide_layouts:
                if tl.name == layout.name:
                    target_layout = tl
                    break
            if target_layout is None:
                target_layout = result_prs.slide_layouts[0]

            new_slide = result_prs.slides.add_slide(target_layout)

            for shape in src_slide.shapes:
                el = deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            # テキストを差し替え
            ph_values = slide_config.get("placeholders", {})
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_name = shape.name
                if shape_name in ph_values:
                    new_text = ph_values[shape_name]["value"]
                    if new_text:
                        _replace_text_preserve_format(shape.text_frame, new_text)
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = ""

        result_prs.save(str(output_path))
        st.session_state.gen_output_path = str(output_path)

        if check_libreoffice():
            try:
                images = render_slides(output_path)
                st.session_state.gen_preview_images = images
            except Exception:
                st.session_state.gen_preview_images = None

    st.rerun()


def _replace_text_preserve_format(text_frame, new_text: str):
    """テキストフレームのテキストを差し替え（書式は保持）。"""
    lines = new_text.split("\n")

    for i, paragraph in enumerate(text_frame.paragraphs):
        if i >= len(lines):
            for run in paragraph.runs:
                run.text = ""
            continue

        if paragraph.runs:
            paragraph.runs[0].text = lines[i]
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = lines[i]

    if len(lines) > len(text_frame.paragraphs):
        for line in lines[len(text_frame.paragraphs):]:
            p = text_frame.add_paragraph()
            p.text = line


# ============================================================
#  チャットモード
# ============================================================

def _render_chat_mode():
    """チャットでコンテンツを伝えてスライドを生成するモード。"""
    st.subheader("チャットでスライドを生成")
    st.caption("作りたいスライドの内容を伝えてください（例: Q1の営業報告を5枚で作って）")

    for msg in st.session_state.gen_chat_messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("スライドの内容を入力"):
        st.session_state.gen_chat_messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            with st.spinner("スライドを生成中..."):
                response = _process_chat_generation(prompt)
                st.markdown(response)

        st.session_state.gen_chat_messages.append({"role": "assistant", "content": response})
        st.rerun()

    # プレビュー
    if st.session_state.gen_preview_images:
        st.divider()
        _show_preview()


def _process_chat_generation(prompt: str) -> str:
    """チャットでの生成指示を処理する。"""
    import anthropic

    client = anthropic.Anthropic()

    template_info = analyze_pptx_summary(st.session_state.gen_template_path)

    # デザイン定義ファイルがあれば読み込む
    template_name = Path(st.session_state.gen_template_path).stem
    designs_dir = Path(__file__).resolve().parent.parent.parent / "designs"
    design_path = designs_dir / f"{template_name}.json"
    design_note = ""
    if design_path.exists():
        with open(design_path, "r", encoding="utf-8") as f:
            design_data = json.load(f)
        design_note = f"\n\nデザイン定義:\n```json\n{json.dumps(design_data, ensure_ascii=False, indent=2)}\n```"

    system_prompt = f"""あなたはスライド生成アシスタントです。
ユーザーの指示に基づいて、テンプレートのスライドにコンテンツを流し込むコードを生成してください。

テンプレート駆動の原則:
- コードでデザインしない（座標・フォント指定禁止）
- テンプレート内のスライドをコピーし、プレースホルダーのテキストを差し替えるだけ
- テンプレートの書式をそのまま継承する

テンプレート構造:
```json
{json.dumps(template_info, ensure_ascii=False, indent=2)}
```
{design_note}

レスポンスは以下のJSON形式で返してください:
```json
{{
  "message": "ユーザーへの返答（生成したスライドの説明）",
  "code": "python-pptxコード"
}}
```

変数:
- `template_path`: テンプレートファイルのパス（文字列）
- `output_path`: 出力先パス（文字列）
- コードの最後で prs.save(output_path) を呼んでください

利用可能なインポート:
- from pptx import Presentation
- from copy import deepcopy
- from pptx.oxml.ns import qn"""

    # チャット履歴を含める
    messages = []
    for msg in st.session_state.gen_chat_messages:
        messages.append({"role": msg["role"], "content": msg["content"]})
    messages.append({"role": "user", "content": prompt})

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        system=system_prompt,
        messages=messages,
    )

    response_text = response.content[0].text

    try:
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        else:
            json_str = response_text
        result = json.loads(json_str)
        message = result.get("message", "スライドを生成しました。")
        code = result.get("code")
    except (json.JSONDecodeError, IndexError):
        return response_text

    if code:
        output_path = _execute_generation_code(code)
        if output_path:
            st.session_state.gen_output_path = output_path
            if check_libreoffice():
                try:
                    images = render_slides(output_path)
                    st.session_state.gen_preview_images = images
                except Exception:
                    pass

    return message


def _execute_generation_code(code: str) -> str | None:
    """生成コードを実行する。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Mm, Emu
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from copy import deepcopy

    template_path = st.session_state.gen_template_path
    output_path = str(Path(tempfile.mkdtemp(prefix="gen_chat_")) / "generated.pptx")

    namespace = {
        "template_path": template_path,
        "output_path": output_path,
        "Presentation": Presentation,
        "Inches": Inches, "Pt": Pt, "Mm": Mm, "Emu": Emu,
        "RGBColor": RGBColor,
        "qn": qn,
        "deepcopy": deepcopy,
    }

    try:
        exec(code, namespace)
    except Exception as e:
        st.error(f"生成コード実行エラー: {e}")
        return None

    if Path(output_path).exists():
        return output_path
    return None


def _show_preview():
    """生成結果のプレビュー。"""
    st.subheader("生成結果プレビュー")
    images = st.session_state.gen_preview_images
    if not images:
        return

    total = len(images)
    if total > 1:
        idx = st.slider("スライド", 0, total - 1, 0, key="gen_preview_idx")
    else:
        idx = 0
    st.image(str(images[idx]), use_container_width=True)
