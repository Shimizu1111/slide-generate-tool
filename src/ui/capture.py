"""デザインキャプチャページ。

画像 / .pptx / PDF をアップロードしてデザインを解析し、
再現可能なテンプレート(.pptx)を作成する。

- 左ペイン: アップロードした元デザインのプレビュー
- 右ペイン: 再現中テンプレートのライブプレビュー
- 下部: チャットで調整指示
"""

from __future__ import annotations

import json
import shutil
import tempfile
from pathlib import Path

import streamlit as st
from PIL import Image

from src.core.analyzer import analyze, analyze_image, analyze_pptx_summary
from src.core.builder import (
    create_blank_presentation,
    add_slide,
    add_textbox,
    add_shape,
    set_slide_background,
    save_presentation,
)
from src.core.renderer import render_slides, check_libreoffice


_SUPPORTED_TYPES = ["pptx", "pdf", "png", "jpg", "jpeg"]


def render_page():
    st.title("デザインキャプチャ")
    st.caption("画像・pptx・PDFをアップロードすると、そのデザインを再現したテンプレートを作ります")

    # セッション状態の初期化
    defaults = {
        "cap_source_path": None,
        "cap_source_type": None,
        "cap_source_preview": None,
        "cap_design_data": None,
        "cap_template_path": None,
        "cap_template_images": [],
        "cap_messages": [],
    }
    for key, default in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default

    # ── サイドバー: ファイルアップロード ──
    with st.sidebar:
        st.subheader("デザインソース")
        source_file = st.file_uploader(
            "ファイルをアップロード",
            type=_SUPPORTED_TYPES,
            key="cap_uploader",
            help="画像 (PNG/JPG)、PowerPoint (.pptx)、PDF に対応",
        )

        if source_file:
            suffix = Path(source_file.name).suffix.lower()
            source_dir = Path(tempfile.mkdtemp(prefix="cap_src_"))
            source_path = source_dir / source_file.name
            source_path.write_bytes(source_file.getvalue())
            st.session_state.cap_source_path = str(source_path)
            st.session_state.cap_source_type = suffix

            # 画像プレビューの準備
            if suffix in {".png", ".jpg", ".jpeg"}:
                st.session_state.cap_source_preview = str(source_path)
            elif suffix == ".pptx" and check_libreoffice():
                try:
                    images = render_slides(source_path, dpi=100)
                    st.session_state.cap_source_preview = str(images[0]) if images else None
                except Exception:
                    st.session_state.cap_source_preview = None
            elif suffix == ".pdf":
                try:
                    from pdf2image import convert_from_path
                    pil_images = convert_from_path(str(source_path), first_page=1, last_page=1, dpi=150)
                    if pil_images:
                        preview_path = source_dir / "preview.png"
                        pil_images[0].save(str(preview_path), "PNG")
                        st.session_state.cap_source_preview = str(preview_path)
                except Exception:
                    st.session_state.cap_source_preview = None

        if st.session_state.cap_source_path and not st.session_state.cap_template_path:
            if st.button("デザインを解析してテンプレート作成", type="primary"):
                _capture_design()

        # テンプレート保存
        if st.session_state.cap_template_path:
            pptx_path = Path(st.session_state.cap_template_path)
            if pptx_path.exists():
                st.divider()
                st.subheader("テンプレートを保存")
                template_name = st.text_input("テンプレート名", value="my_template", key="cap_template_name")

                if st.button("templates/ に保存"):
                    templates_dir = Path(__file__).resolve().parent.parent.parent / "templates"
                    templates_dir.mkdir(exist_ok=True)
                    dest = templates_dir / f"{template_name}.pptx"
                    shutil.copy2(str(pptx_path), str(dest))

                    # デザイン定義も保存
                    if st.session_state.cap_design_data:
                        designs_dir = Path(__file__).resolve().parent.parent.parent / "designs"
                        designs_dir.mkdir(exist_ok=True)
                        design_path = designs_dir / f"{template_name}.json"
                        with open(design_path, "w", encoding="utf-8") as f:
                            json.dump(st.session_state.cap_design_data, f, ensure_ascii=False, indent=2)
                        st.success(f"保存: templates/{template_name}.pptx + designs/{template_name}.json")
                    else:
                        st.success(f"保存: templates/{template_name}.pptx")

                with open(pptx_path, "rb") as f:
                    st.download_button(
                        label="ダウンロード",
                        data=f.read(),
                        file_name=f"{template_name}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

    # ── メインエリア ──
    if not st.session_state.cap_source_path:
        st.info("サイドバーからファイルをアップロードしてください")
        return

    # 左右比較
    left_col, right_col = st.columns([1, 1])

    with left_col:
        st.subheader("元のデザイン")
        if st.session_state.cap_source_preview:
            st.image(st.session_state.cap_source_preview, use_container_width=True)
        else:
            st.info("プレビューを表示できません")

    with right_col:
        st.subheader("再現テンプレート")
        if st.session_state.cap_template_images:
            images = st.session_state.cap_template_images
            total = len(images)
            if total > 1:
                idx = st.slider("スライド", 0, total - 1, 0, key="cap_slide_idx")
            else:
                idx = 0
            st.image(str(images[idx]), use_container_width=True)
        else:
            st.info("「デザインを解析してテンプレート作成」を押してください")

    # ── デザイン定義の表示 ──
    if st.session_state.cap_design_data:
        with st.expander("解析されたデザイン定義"):
            st.json(st.session_state.cap_design_data)

    # ── 調整チャット ──
    if st.session_state.cap_template_path:
        st.divider()
        _show_refinement_chat()


def _capture_design():
    """デザインを解析してテンプレートを生成する。"""
    source_path = st.session_state.cap_source_path
    source_type = st.session_state.cap_source_type

    with st.spinner("デザインを解析中..."):
        try:
            if source_type == ".pptx":
                # pptxの場合: そのままコピーしてテンプレート化
                design_data = analyze_pptx_summary(source_path)
                template_path = _clone_pptx_as_template(source_path)
            else:
                # 画像/PDFの場合: Vision APIで解析→builder で構築
                design_data = analyze(source_path)
                template_path = _build_template_from_design(design_data)

            st.session_state.cap_design_data = design_data
            st.session_state.cap_template_path = template_path
            _update_template_preview()
            st.rerun()

        except Exception as e:
            st.error(f"デザイン解析エラー: {e}")


def _clone_pptx_as_template(source_path: str) -> str:
    """pptxをコピーしてプレースホルダーのテキストをクリアする。"""
    from pptx import Presentation

    output_dir = Path(tempfile.mkdtemp(prefix="cap_tpl_"))
    output_path = output_dir / "template.pptx"
    shutil.copy2(source_path, str(output_path))

    prs = Presentation(str(output_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                ph = shape.placeholder_format
                if ph is not None and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
            except ValueError:
                continue

    prs.save(str(output_path))
    return str(output_path)


def _build_template_from_design(design_data: dict) -> str:
    """デザイン解析結果からテンプレートを構築する。"""
    width = design_data.get("slide_width_mm", 254.0)
    height = design_data.get("slide_height_mm", 143.0)

    prs = create_blank_presentation(width, height)
    slide = add_slide(prs, 0)

    # 背景
    bg = design_data.get("background", {})
    if bg.get("type") == "solid" and bg.get("color"):
        set_slide_background(slide, bg["color"])

    # 要素を配置
    for elem in design_data.get("elements", []):
        elem_type = elem.get("type", "")

        if elem_type == "textbox":
            add_textbox(
                slide,
                left_mm=elem.get("left_mm", 0),
                top_mm=elem.get("top_mm", 0),
                width_mm=elem.get("width_mm", 100),
                height_mm=elem.get("height_mm", 30),
                text=elem.get("text", ""),
                font_name=elem.get("font_name"),
                font_size_pt=elem.get("font_size_pt"),
                font_color_rgb=elem.get("font_color"),
                bold=elem.get("bold"),
                italic=elem.get("italic"),
                alignment=elem.get("alignment"),
            )

        elif elem_type == "shape":
            add_shape(
                slide,
                shape_type=elem.get("shape_type", "rectangle"),
                left_mm=elem.get("left_mm", 0),
                top_mm=elem.get("top_mm", 0),
                width_mm=elem.get("width_mm", 100),
                height_mm=elem.get("height_mm", 30),
                fill_color_rgb=elem.get("fill_color"),
                line_color_rgb=elem.get("line_color"),
            )

    output_path = Path(tempfile.mkdtemp(prefix="cap_tpl_")) / "template.pptx"
    save_presentation(prs, output_path)
    return str(output_path)


def _show_refinement_chat():
    """テンプレート調整用チャット。"""
    st.subheader("テンプレートを調整")
    st.caption("元のデザインと見比べて、違いがあればチャットで指示してください")

    for msg in st.session_state.cap_messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("調整指示を入力（例: タイトルのフォントをもう少し大きく）"):
        st.session_state.cap_messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            with st.spinner("修正中..."):
                response = _process_refinement(prompt)
                st.markdown(response)

        st.session_state.cap_messages.append({"role": "assistant", "content": response})
        st.rerun()


def _process_refinement(prompt: str) -> str:
    """調整指示を処理する。"""
    import anthropic

    client = anthropic.Anthropic()

    template_info = analyze_pptx_summary(st.session_state.cap_template_path)

    # 元デザインの画像があれば添付
    content_parts = []
    source_preview = st.session_state.cap_source_preview
    if source_preview and Path(source_preview).exists():
        import base64
        image_data = base64.standard_b64encode(Path(source_preview).read_bytes()).decode("utf-8")
        suffix = Path(source_preview).suffix.lower()
        media_type = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg"}.get(
            suffix.lstrip("."), "image/png"
        )
        content_parts.append({
            "type": "image",
            "source": {"type": "base64", "media_type": media_type, "data": image_data},
        })

    # テンプレートのプレビュー画像も添付
    template_images = st.session_state.cap_template_images
    if template_images:
        import base64
        first_img = Path(template_images[0])
        if first_img.exists():
            img_data = base64.standard_b64encode(first_img.read_bytes()).decode("utf-8")
            content_parts.append({
                "type": "image",
                "source": {"type": "base64", "media_type": "image/png", "data": img_data},
            })

    content_parts.append({
        "type": "text",
        "text": (
            f"現在のテンプレート構造:\n```json\n{json.dumps(template_info, ensure_ascii=False, indent=2)}\n```\n\n"
            f"1枚目の画像が元のデザイン、2枚目が現在のテンプレートです。\n"
            f"ユーザーの指示: {prompt}"
        ),
    })

    system_prompt = """あなたはpptxテンプレートを調整するエキスパートです。
ユーザーの指示に基づいて、テンプレートを修正するpython-pptxコードを生成してください。

レスポンスは以下のJSON形式で返してください:
```json
{
  "message": "ユーザーへの返答",
  "code": "修正用Pythonコード"
}
```

変数 `prs` に現在のPresentationオブジェクトが入っています。
修正後も `prs` に残してください。save()は呼ばないでください。

利用可能なインポート:
- from pptx import Presentation
- from pptx.util import Inches, Pt, Mm, Emu
- from pptx.dml.color import RGBColor
- from pptx.enum.text import PP_ALIGN, MSO_ANCHOR"""

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        system=system_prompt,
        messages=[{"role": "user", "content": content_parts}],
    )

    response_text = response.content[0].text

    try:
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        else:
            json_str = response_text
        result = json.loads(json_str)
        message = result.get("message", "修正しました。")
        code = result.get("code")
    except (json.JSONDecodeError, IndexError):
        return response_text

    if code:
        new_path = _execute_refinement_code(code)
        if new_path:
            st.session_state.cap_template_path = new_path
            _update_template_preview()

    return message


def _execute_refinement_code(code: str) -> str | None:
    """修正コードを実行する。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Mm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation(st.session_state.cap_template_path)

    namespace = {
        "prs": prs,
        "Presentation": Presentation,
        "Inches": Inches, "Pt": Pt, "Mm": Mm, "Emu": Emu,
        "RGBColor": RGBColor,
        "PP_ALIGN": PP_ALIGN, "MSO_ANCHOR": MSO_ANCHOR,
    }

    try:
        exec(code, namespace)
    except Exception as e:
        st.error(f"修正コード実行エラー: {e}")
        return None

    prs = namespace.get("prs")
    if prs is None:
        return None

    output_path = Path(tempfile.mkdtemp(prefix="cap_refined_")) / "template.pptx"
    save_presentation(prs, output_path)
    return str(output_path)


def _update_template_preview():
    """テンプレートのプレビュー画像を更新する。"""
    path = st.session_state.cap_template_path
    if not path or not Path(path).exists():
        st.session_state.cap_template_images = []
        return

    if not check_libreoffice():
        st.session_state.cap_template_images = []
        return

    try:
        images = render_slides(path)
        st.session_state.cap_template_images = images
    except Exception:
        st.session_state.cap_template_images = []
