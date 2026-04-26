"""テンプレート作成ページ。

左ペイン: チャットでデザイン意図を伝える
右ペイン: テンプレートのライブプレビュー
"""

from __future__ import annotations

import json
import tempfile
from pathlib import Path

import streamlit as st

from src.core.builder import (
    create_blank_presentation,
    add_slide,
    add_textbox,
    add_shape,
    set_slide_background,
    save_presentation,
)
from src.core.renderer import render_slides, check_libreoffice
from src.core.inspector import inspect_template


def render_page():
    st.title("テンプレート作成")
    st.caption("チャットでデザイン意図を伝え、テンプレートを作成します")

    # セッション状態の初期化
    if "create_messages" not in st.session_state:
        st.session_state.create_messages = []
    if "create_pptx_path" not in st.session_state:
        st.session_state.create_pptx_path = None
    if "create_preview_images" not in st.session_state:
        st.session_state.create_preview_images = []

    # 2カラムレイアウト
    chat_col, preview_col = st.columns([1, 1])

    # ── 左ペイン: チャット ──
    with chat_col:
        st.subheader("チャット")

        # メッセージ履歴の表示
        for msg in st.session_state.create_messages:
            with st.chat_message(msg["role"]):
                st.markdown(msg["content"])

        # ユーザー入力
        if prompt := st.chat_input("デザインの要望を入力してください"):
            st.session_state.create_messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)

            # AIレスポンス生成
            with st.chat_message("assistant"):
                with st.spinner("テンプレートを生成中..."):
                    response, pptx_path = _process_chat(
                        st.session_state.create_messages,
                        st.session_state.create_pptx_path,
                    )
                    st.markdown(response)

            st.session_state.create_messages.append({"role": "assistant", "content": response})
            if pptx_path:
                st.session_state.create_pptx_path = pptx_path
                _update_preview()
            st.rerun()

    # ── 右ペイン: プレビュー ──
    with preview_col:
        st.subheader("プレビュー")

        if st.session_state.create_preview_images:
            # スライド番号選択
            total = len(st.session_state.create_preview_images)
            if total > 1:
                slide_idx = st.slider("スライド", 0, total - 1, 0, key="create_slide_idx")
            else:
                slide_idx = 0

            st.image(
                str(st.session_state.create_preview_images[slide_idx]),
                use_container_width=True,
            )

            # テンプレート情報表示
            if st.session_state.create_pptx_path:
                with st.expander("テンプレート構造"):
                    info = inspect_template(st.session_state.create_pptx_path)
                    st.json(info)
        else:
            st.info("チャットでデザインの要望を伝えると、プレビューが表示されます")

        # 保存・ダウンロード
        if st.session_state.create_pptx_path:
            pptx_path = Path(st.session_state.create_pptx_path)
            if pptx_path.exists():
                template_name = st.text_input("テンプレート名", value="my_template", key="create_template_name")

                col_save, col_dl = st.columns(2)
                with col_save:
                    if st.button("templates/ に保存"):
                        import shutil
                        templates_dir = Path(__file__).resolve().parent.parent.parent / "templates"
                        templates_dir.mkdir(exist_ok=True)
                        dest = templates_dir / f"{template_name}.pptx"
                        shutil.copy2(str(pptx_path), str(dest))
                        st.success(f"保存: templates/{template_name}.pptx")
                        st.code(f"make generate T=templates/{template_name}.pptx I=input/data.json", language="bash")
                with col_dl:
                    with open(pptx_path, "rb") as f:
                        st.download_button(
                            label="ダウンロード",
                            data=f.read(),
                            file_name=f"{template_name}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )


def _process_chat(
    messages: list[dict],
    current_pptx_path: str | None,
) -> tuple[str, str | None]:
    """チャットメッセージを処理し、テンプレートを生成・更新する。

    OpenAI APIを使ってユーザーの意図を解釈し、
    builder関数群を呼び出してテンプレートを構築する。
    """
    from openai import OpenAI

    client = OpenAI()

    # 現在のテンプレート情報をコンテキストに含める
    context = ""
    if current_pptx_path and Path(current_pptx_path).exists():
        info = inspect_template(current_pptx_path)
        context = f"\n\n現在のテンプレートの構造:\n```json\n{json.dumps(info, ensure_ascii=False, indent=2)}\n```"

    system_prompt = f"""あなたはスライドテンプレートを設計するアシスタントです。
ユーザーの要望に基づいてpython-pptxのコードを生成し、テンプレートを構築します。

以下のbuilder関数が利用可能です:
- create_blank_presentation(width_mm, height_mm) → Presentation
- add_slide(prs, layout_index, content) → slide
- add_textbox(slide, left_mm, top_mm, width_mm, height_mm, text, font_name, font_size_pt, font_color_rgb, bold, italic, alignment, vertical_anchor)
- add_shape(slide, shape_type, left_mm, top_mm, width_mm, height_mm, fill_color_rgb, line_color_rgb, line_width_pt)
- set_slide_background(slide, color_rgb)
- save_presentation(prs, output_path)

shape_typeの選択肢: rectangle, rounded_rectangle, oval, triangle, arrow_right, line
alignmentの選択肢: left, center, right, justify
vertical_anchorの選択肢: top, middle, bottom
color_rgbの例: "FF5733", "333333", "FFFFFF"

レスポンスは以下のJSON形式で返してください:
```json
{{
  "message": "ユーザーへの返答テキスト",
  "code": "実行するPythonコード（builder関数を使用）"
}}
```

codeフィールドには、以下のインポートが既に行われている前提でコードを書いてください:
- from src.core.builder import *
- 変数 `prs` に現在のPresentationオブジェクトが入っています（新規の場合はNone）
- コードの最後で `prs` 変数にPresentationオブジェクトを代入してください
- save_presentationは呼ばないでください（自動的に呼ばれます）

コードがない場合（質問への回答のみの場合）はcodeフィールドをnullにしてください。
{context}"""

    # API呼び出し用メッセージ構築
    api_messages = []
    for msg in messages:
        api_messages.append({"role": msg["role"], "content": msg["content"]})

    response = client.chat.completions.create(
        model="gpt-5.4",
        max_completion_tokens=4096,
        messages=[{"role": "system", "content": system_prompt}] + api_messages,
    )

    response_text = response.choices[0].message.content

    # JSONレスポンスを解析
    try:
        # コードブロック内のJSONを抽出
        if "```json" in response_text:
            json_str = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            json_str = response_text.split("```")[1].split("```")[0].strip()
        else:
            json_str = response_text

        result = json.loads(json_str)
        message = result.get("message", "テンプレートを更新しました。")
        code = result.get("code")
    except (json.JSONDecodeError, IndexError):
        return response_text, None

    if not code:
        return message, None

    # コードを実行してテンプレートを生成
    pptx_path = _execute_builder_code(code, current_pptx_path)
    return message, pptx_path


def _execute_builder_code(code: str, current_pptx_path: str | None) -> str | None:
    """builder関数を使ったコードを実行する。"""
    from src.core.builder import (
        create_blank_presentation,
        add_slide,
        add_textbox,
        add_shape,
        add_image,
        set_slide_background,
        save_presentation,
        modify_layout,
    )
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Mm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # 現在のプレゼンテーションを読み込み
    prs = None
    if current_pptx_path and Path(current_pptx_path).exists():
        prs = PptxPresentation(current_pptx_path)

    # コード実行用の名前空間
    namespace = {
        "prs": prs,
        "create_blank_presentation": create_blank_presentation,
        "add_slide": add_slide,
        "add_textbox": add_textbox,
        "add_shape": add_shape,
        "add_image": add_image,
        "set_slide_background": set_slide_background,
        "save_presentation": save_presentation,
        "modify_layout": modify_layout,
        "Presentation": PptxPresentation,
        "Inches": Inches,
        "Pt": Pt,
        "Mm": Mm,
        "Emu": Emu,
        "RGBColor": RGBColor,
        "PP_ALIGN": PP_ALIGN,
        "MSO_ANCHOR": MSO_ANCHOR,
    }

    try:
        exec(code, namespace)
    except Exception as e:
        st.error(f"コード実行エラー: {e}")
        return None

    prs = namespace.get("prs")
    if prs is None:
        st.warning("Presentationオブジェクトが生成されませんでした。")
        return None

    # 一時ファイルに保存
    output_path = Path(tempfile.mkdtemp(prefix="template_")) / "template.pptx"
    save_presentation(prs, output_path)
    return str(output_path)


def _update_preview():
    """プレビュー画像を更新する。"""
    pptx_path = st.session_state.create_pptx_path
    if not pptx_path or not Path(pptx_path).exists():
        st.session_state.create_preview_images = []
        return

    if not check_libreoffice():
        st.warning("LibreOfficeがインストールされていないため、プレビューを表示できません。")
        st.session_state.create_preview_images = []
        return

    try:
        images = render_slides(pptx_path)
        st.session_state.create_preview_images = images
    except Exception as e:
        st.error(f"プレビュー生成エラー: {e}")
        st.session_state.create_preview_images = []
